#!/usr/bin/env python3
"""
Genera presentación de ventas de CASTOR Elecciones para candidatos/campañas.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colores de marca CASTOR
CASTOR_BLUE = RGBColor(0x1E, 0x3A, 0x5F)  # Azul oscuro
CASTOR_GOLD = RGBColor(0xD4, 0xA5, 0x37)  # Dorado
CASTOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CASTOR_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
CASTOR_DARK = RGBColor(0x2C, 0x3E, 0x50)

def set_slide_background(slide, color):
    """Establecer color de fondo del slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle=""):
    """Agregar slide de título."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CASTOR_BLUE)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = CASTOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = CASTOR_GOLD
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight=None):
    """Agregar slide con contenido en bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CASTOR_WHITE)

    # Barra superior azul
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CASTOR_BLUE
    shape.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CASTOR_WHITE

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = CASTOR_DARK
        p.space_after = Pt(12)

        # Resaltar item específico
        if highlight and bullet == highlight:
            p.font.bold = True
            p.font.color.rgb = CASTOR_BLUE

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Agregar slide con dos columnas."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CASTOR_WHITE)

    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CASTOR_BLUE
    shape.line.fill.background()

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CASTOR_WHITE

    # Columna izquierda
    left_header = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CASTOR_BLUE

    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = CASTOR_DARK
        p.space_after = Pt(8)

    # Columna derecha
    right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CASTOR_GOLD

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.1), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = CASTOR_DARK
        p.space_after = Pt(8)

    return slide

def add_stats_slide(prs, title, stats):
    """Agregar slide con estadísticas grandes."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CASTOR_BLUE)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = CASTOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Stats en fila
    num_stats = len(stats)
    width = 9 / num_stats

    for i, (number, label) in enumerate(stats):
        x = 0.5 + (i * width)

        # Número grande
        num_box = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(width), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = CASTOR_GOLD
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(width), Inches(1))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = CASTOR_WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_feature_slide(prs, title, feature_name, description, benefits):
    """Agregar slide de feature con descripción y beneficios."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CASTOR_WHITE)

    # Barra superior
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CASTOR_BLUE
    shape.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CASTOR_WHITE

    # Nombre del feature
    feat_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.6))
    tf = feat_box.text_frame
    p = tf.paragraphs[0]
    p.text = feature_name
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CASTOR_GOLD

    # Descripción
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.2))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(18)
    p.font.color.rgb = CASTOR_DARK

    # Beneficios
    ben_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(3))
    tf = ben_box.text_frame
    tf.word_wrap = True

    for i, benefit in enumerate(benefits):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f">> {benefit}"
        p.font.size = Pt(18)
        p.font.color.rgb = CASTOR_BLUE
        p.space_after = Pt(10)

    return slide

def add_cta_slide(prs, title, subtitle, cta_text, contact_info):
    """Agregar slide de llamado a la acción."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, CASTOR_BLUE)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = CASTOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = CASTOR_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # CTA Button (simulado con shape)
    cta_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(3.5), Inches(4), Inches(0.8))
    cta_shape.fill.solid()
    cta_shape.fill.fore_color.rgb = CASTOR_GOLD
    cta_shape.line.fill.background()

    cta_box = slide.shapes.add_textbox(Inches(3), Inches(3.6), Inches(4), Inches(0.6))
    tf = cta_box.text_frame
    p = tf.paragraphs[0]
    p.text = cta_text
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = CASTOR_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Contacto
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    tf = contact_box.text_frame
    for i, info in enumerate(contact_info):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = info
        p.font.size = Pt(16)
        p.font.color.rgb = CASTOR_LIGHT
        p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Crear la presentación completa."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: PORTADA =====
    add_title_slide(
        prs,
        "CASTOR ELECCIONES",
        "Inteligencia Electoral en Tiempo Real"
    )

    # ===== SLIDE 2: EL PROBLEMA =====
    add_content_slide(
        prs,
        "El Problema: Campanas a Ciegas",
        [
            "Las campanas tradicionales no saben que piensa la gente en tiempo real",
            "Encuestas costosas y desactualizadas (semanas de retraso)",
            "Imposible medir el impacto de cada mensaje o evento",
            "Los rivales se mueven y no hay forma de reaccionar rapido",
            "Decisiones basadas en intuicion, no en datos",
            "Miles de conversaciones en redes sociales que nunca se analizan"
        ]
    )

    # ===== SLIDE 3: LA SOLUCIÓN =====
    add_content_slide(
        prs,
        "La Solucion: CASTOR Elecciones",
        [
            "Analisis de sentimiento en tiempo real de Twitter/X",
            "Inteligencia artificial entrenada en espanol colombiano",
            "Dashboard con metricas electorales exclusivas",
            "Predicciones basadas en conversacion digital",
            "Alertas de cambios de tendencia",
            "Asistente IA para estrategia de campana"
        ]
    )

    # ===== SLIDE 4: CÓMO FUNCIONA =====
    add_two_column_slide(
        prs,
        "Como Funciona CASTOR?",
        "Recoleccion",
        [
            "Captura tweets en tiempo real",
            "Filtra por ubicacion geografica",
            "Detecta menciones de candidatos",
            "Identifica temas trending",
            "Procesa miles de tweets por hora"
        ],
        "Analisis",
        [
            "IA analiza sentimiento (BETO)",
            "Clasifica por tema politico",
            "Calcula metricas exclusivas",
            "Genera insights automaticos",
            "Entrega resumenes ejecutivos"
        ]
    )

    # ===== SLIDE 5: ESTADÍSTICAS =====
    add_stats_slide(
        prs,
        "CASTOR en Numeros",
        [
            ("10K+", "Tweets/dia\nanalizados"),
            ("95%", "Precision en\nsentimiento"),
            ("< 5 min", "Tiempo de\nactualizacion"),
            ("32", "Departamentos\nde Colombia")
        ]
    )

    # ===== SLIDE 6: ICCE =====
    add_feature_slide(
        prs,
        "Metricas Exclusivas",
        "ICCE - Indice Compuesto de Conversacion Electoral",
        "Un numero unico (0-100) que resume la salud de tu campana en la conversacion digital. Combina volumen, sentimiento, y engagement en una metrica facil de entender.",
        [
            "Compara tu ICCE con el de tus rivales",
            "Monitorea la evolucion dia a dia",
            "Identifica que acciones mejoran tu indice",
            "Detecta crisis antes de que escalen"
        ]
    )

    # ===== SLIDE 7: MOMENTUM =====
    add_feature_slide(
        prs,
        "Metricas Exclusivas",
        "Momentum Electoral",
        "Mide la velocidad de cambio de tu campana. Estas ganando terreno o perdiendolo? El momentum te dice si vas en la direccion correcta.",
        [
            "Detecta tendencias antes que las encuestas",
            "Identifica el efecto de debates y eventos",
            "Compara momentum entre candidatos",
            "Anticipa cambios en intencion de voto"
        ]
    )

    # ===== SLIDE 8: PREDICCIONES =====
    add_feature_slide(
        prs,
        "Predicciones Electorales",
        "Forecast con Inteligencia Artificial",
        "Proyectamos el ICCE a 14 dias con intervalos de confianza. Simulamos escenarios para que planifiques tu estrategia.",
        [
            "Proyecciones a 7, 14 y 30 dias",
            "Simulacion de escenarios (debates, escandalos, endorsements)",
            "Intervalos de confianza estadisticos",
            "Alertas automaticas de cambios significativos"
        ]
    )

    # ===== SLIDE 9: ANÁLISIS POR TEMA =====
    add_content_slide(
        prs,
        "Analisis por Tema del PND",
        [
            "Seguridad ciudadana y defensa nacional",
            "Economia, empleo y desarrollo",
            "Salud y bienestar",
            "Educacion y cultura",
            "Medio ambiente y cambio climatico",
            "Infraestructura y transporte",
            "Paz y reconciliacion",
            "Corrupcion y transparencia"
        ]
    )

    # ===== SLIDE 10: RIVALES =====
    add_feature_slide(
        prs,
        "Inteligencia Competitiva",
        "Analisis de Rivales",
        "Monitorea en tiempo real lo que dicen de tus competidores. Identifica sus fortalezas, debilidades y oportunidades para diferenciarte.",
        [
            "Comparativa de ICCE entre candidatos",
            "Analisis de narrativas de cada rival",
            "Deteccion de ataques y respuestas",
            "Oportunidades de posicionamiento"
        ]
    )

    # ===== SLIDE 11: ASISTENTE IA =====
    add_feature_slide(
        prs,
        "Asistente de Campana",
        "Chat IA con Memoria",
        "Un asistente inteligente que conoce tu campana, tus datos y te ayuda a tomar decisiones. Preguntale cualquier cosa sobre tu estrategia.",
        [
            "Respuestas basadas en TUS datos de campana",
            "Sugerencias de mensajes y contenido",
            "Analisis de que funciono y que no",
            "Disponible 24/7 para tu equipo"
        ]
    )

    # ===== SLIDE 12: CASO DE USO =====
    add_two_column_slide(
        prs,
        "Caso de Uso: Dia de Debate",
        "Sin CASTOR",
        [
            "Ves el debate y esperas",
            "Revisas redes manualmente",
            "Encuesta en 2 semanas",
            "No sabes si ganaste o perdiste",
            "Reaccionas tarde a la narrativa"
        ],
        "Con CASTOR",
        [
            "ICCE en tiempo real durante debate",
            "Alertas de picos de sentimiento",
            "Sabes en minutos quien gano",
            "Identificas temas que resonaron",
            "Ajustas narrativa inmediatamente"
        ]
    )

    # ===== SLIDE 13: BENEFICIOS =====
    add_content_slide(
        prs,
        "Beneficios Clave para tu Campana",
        [
            "Toma decisiones basadas en datos, no en intuicion",
            "Reacciona en horas, no en semanas",
            "Optimiza tu presupuesto de comunicacion",
            "Anticipa crisis antes de que escalen",
            "Entiende que mensajes resuenan con los votantes",
            "Gana la batalla de la narrativa digital"
        ]
    )

    # ===== SLIDE 14: PLANES =====
    add_two_column_slide(
        prs,
        "Planes y Precios",
        "Plan Campana",
        [
            "Dashboard completo",
            "ICCE y Momentum",
            "Analisis de 1 rival",
            "Asistente IA basico",
            "Actualizaciones cada hora",
            "Soporte por email",
            "",
            "Desde $2.5M COP/mes"
        ],
        "Plan Estratega",
        [
            "Todo del Plan Campana +",
            "Analisis de 5 rivales",
            "Forecast y simulaciones",
            "Asistente IA avanzado",
            "Actualizaciones tiempo real",
            "Soporte prioritario 24/7",
            "Reportes personalizados",
            "Desde $7M COP/mes"
        ]
    )

    # ===== SLIDE 15: POR QUÉ CASTOR =====
    add_content_slide(
        prs,
        "Por Que Elegir CASTOR?",
        [
            "Tecnologia colombiana, para campanas colombianas",
            "IA entrenada en espanol y contexto politico local",
            "Equipo con experiencia en campanas reales",
            "Metricas exclusivas que nadie mas ofrece",
            "Soporte dedicado durante toda la campana",
            "Confidencialidad total de tus datos"
        ]
    )

    # ===== SLIDE 16: TECNOLOGÍA =====
    add_content_slide(
        prs,
        "Tecnologia de Clase Mundial",
        [
            "BETO: Modelo de IA entrenado en espanol",
            "GPT-4: Generacion de insights y recomendaciones",
            "ChromaDB: Base de conocimiento con memoria",
            "Twitter/X API: Datos en tiempo real",
            "Infraestructura en la nube, siempre disponible",
            "Encriptacion de extremo a extremo"
        ]
    )

    # ===== SLIDE 17: TESTIMONIAL =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, CASTOR_LIGHT)

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"CASTOR nos permitio detectar una crisis de reputacion 48 horas antes de que explotara en medios tradicionales. Pudimos reaccionar a tiempo y convertir un potencial desastre en una oportunidad."'
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = CASTOR_DARK
    p.alignment = PP_ALIGN.CENTER

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    tf = attr_box.text_frame
    p = tf.paragraphs[0]
    p.text = "- Director de Campana, Elecciones Regionales 2023"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CASTOR_BLUE
    p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 18: PRÓXIMOS PASOS =====
    add_content_slide(
        prs,
        "Proximos Pasos",
        [
            "1. Agenda una demo personalizada (30 min)",
            "2. Te mostramos datos reales de tu region",
            "3. Disenamos un plan a la medida de tu campana",
            "4. Activamos tu cuenta en 24 horas",
            "5. Capacitamos a tu equipo",
            "6. Empiezas a ganar con datos!"
        ]
    )

    # ===== SLIDE 19: CTA FINAL =====
    add_cta_slide(
        prs,
        "Listo para Ganar con Inteligencia?",
        "Las elecciones se ganan con informacion. Nosotros la tenemos.",
        "SOLICITAR DEMO GRATIS",
        [
            "contacto@castorelecciones.com",
            "+57 300 123 4567",
            "www.castorelecciones.com"
        ]
    )

    # ===== SLIDE 20: GRACIAS =====
    add_title_slide(
        prs,
        "CASTOR ELECCIONES",
        "Inteligencia que Gana Elecciones"
    )

    return prs


if __name__ == "__main__":
    print("Generando presentacion de ventas...")
    prs = create_presentation()

    output_path = "/Users/arielsanroj/castor/CASTOR_Ventas_Campanas.pptx"
    prs.save(output_path)
    print(f"Presentacion guardada en: {output_path}")
    print(f"Total de slides: {len(prs.slides)}")
